import pyautogui
from PIL import Image
from pynput import mouse
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import ttk, simpledialog, messagebox
import threading
import os
import webbrowser
import subprocess
import platform


# Global variables
recording = False
screenshot_index = 0
captured_images = []
folder_name = ""

def capture_screenshot():
    global screenshot_index
    if not folder_name:
        messagebox.showwarning("Folder Error", "Folder not set. Please start recording again.")
        return

    screenshot = pyautogui.screenshot()
    filename = os.path.join(folder_name, f'step_{screenshot_index}.png')
    screenshot.save(filename)
    captured_images.append(filename)
    screenshot_index += 1

def on_click(x, y, button, pressed):
    if pressed and recording:
        capture_screenshot()

def start_recording():
    global recording
    if not folder_name:
        messagebox.showwarning("Folder Error", "Folder not created. Please create a folder first.")
        return
    recording = True
    messagebox.showinfo("Recording", "Recording started. Click to capture screenshots.")

def stop_recording():
    global recording
    recording = False

    if folder_name:
        pptx_name = simpledialog.askstring("Input", "Enter the name for the PowerPoint presentation (without extension):")
        if pptx_name:
            save_to_pptx(pptx_name)
            messagebox.showinfo("Recording Stopped", f"Recording stopped. Files saved in '{folder_name}' folder.")
        else:
            messagebox.showwarning("No Name Provided", "PowerPoint presentation name was not provided.")
    else:
        messagebox.showwarning("No Folder Created", "Folder was not created. Please start again.")

def save_to_pptx(pptx_name):
    prs = Presentation()
    for image in captured_images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Empty slide layout
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image, left, top, height=Inches(5.5))
    prs.save(os.path.join(folder_name, f'{pptx_name}.pptx'))

def show_directions():
    directions_text = (
        "Directions to the Step Recorder\n"
        "1) Click the Create Folder button to name the folder the images and the PowerPoint (pptx file) will be exported to once the recording has finished. This folder will default to the desktop.\n\n"
        "2) Click Start Recording to start recording your steps. An image will be created with each mouse click.\n\n"
        "3) Click Stop Recording to stop the recording of your steps.\n\n"
        "4) Open the folder with your recording. You will have images from each mouse click. You will have a PowerPoint presentation which you will be able to edit. Name the step, delete steps if needed and use any transitions and effects within PowerPoint. Please save once you are finished. You can also export the pptx as a video file.\n\n"
	"5) A folder will be created with a powerpoint of each image on a slide; where you can edit the powerpoint; adding captaion, shapes, text, voice recording, links to specific sites, etc. You will also have each individual photo image that was created with each mouse click.\n\n"
        "6) ----ADVANCE - EXPORTING----\n"
        "Export PowerPoint (pptx) as webpage (html) at: PPTX2HTML below"
    )

    directions_popup = tk.Toplevel()
    directions_popup.title("Directions")

    text = tk.Text(directions_popup, wrap=tk.WORD, height=25, width=80)
    text.insert(tk.END, directions_text)
    text.config(state=tk.DISABLED)
    text.pack()

    link = tk.Label(directions_popup, text="PPTX2HTML", fg="blue", cursor="hand2")
    link.pack(pady=10)
    link.bind("<Button-1>", lambda e: webbrowser.open_new("https://tools.gokhs.org/pptx2html"))

    ok_button = ttk.Button(directions_popup, text="OK", command=directions_popup.destroy)
    ok_button.pack(pady=10)

def create_gui():
    root = tk.Tk()
    root.title("Step Recorder")
    
    root.iconbitmap("KSR.ico")

    header_label = tk.Label(root, text="Steps Recorder", font=("Helvetica", 16))
    header_label.pack(pady=10)

    directions_button = ttk.Button(root, text="Directions", command=show_directions)
    directions_button.pack(pady=10)

    create_folder_button = ttk.Button(root, text="Create Folder", command=create_folder)
    create_folder_button.pack(pady=10)

    start_button = ttk.Button(root, text="Start Recording", command=start_recording)
    start_button.pack(pady=10)

    stop_button = ttk.Button(root, text="Stop Recording", command=stop_recording)
    stop_button.pack(pady=10)

    quit_button = ttk.Button(root, text="Quit", command=root.quit)
    quit_button.pack(pady=10)

    root.mainloop()

def create_folder():
    global folder_name
    folder_name = os.path.join(os.path.expanduser("~"), "Desktop", simpledialog.askstring("Input", "Enter the name for the folder:"))
    if folder_name:
        if not os.path.exists(folder_name):
            os.makedirs(folder_name)
        messagebox.showinfo("Folder Created", f"Folder '{folder_name}' created. You can now start recording.")
    else:
        messagebox.showwarning("No Folder Name Provided", "Folder name was not provided.")

def start_mouse_listener():
    """Start listening for mouse clicks."""
    with mouse.Listener(on_click=on_click) as listener:
        listener.join()

if __name__ == "__main__":
    # Run the GUI in one thread and the mouse listener in another
    gui_thread = threading.Thread(target=create_gui)
    mouse_listener_thread = threading.Thread(target=start_mouse_listener)

    gui_thread.start()
    mouse_listener_thread.start()

    gui_thread.join()
    mouse_listener_thread.join()
